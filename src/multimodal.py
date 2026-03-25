"""
src/multimodal.py — Multimodal brief extraction
Converts PDF and PPTX campaign briefs into text + JPEG slide images.
"""
import base64
import io
import os
import subprocess
import tempfile


def _pdf_bytes_to_images(pdf_bytes: bytes, max_slides: int, max_dim: int) -> list[str]:
    """Render each page of a PDF to a JPEG and return list of base64 strings."""
    import fitz  # PyMuPDF
    from PIL import Image

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    images: list[str] = []

    for page_num in range(min(len(doc), max_slides)):
        page = doc[page_num]
        # Render at 1.5× for sharpness, then thumbnail to max_dim
        mat = fitz.Matrix(1.5, 1.5)
        pix = page.get_pixmap(matrix=mat)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        img.thumbnail((max_dim, max_dim), Image.LANCZOS)

        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=75, optimize=True)
        images.append(base64.b64encode(buf.getvalue()).decode())

    doc.close()
    return images


def extract_multimodal_brief(
    file_bytes: bytes,
    file_type: str,
    max_slides: int = 6,
    max_dim: int = 1024,
) -> dict:
    """
    Extract text + slide images from a PDF or PPTX brief.

    Returns:
        {
            "text": str,            # Full extracted text
            "images": [str, ...],   # Base64 JPEG strings, max max_slides
            "slide_count": int,
        }
    """
    file_type = file_type.lower().lstrip(".")
    text = ""
    images: list[str] = []

    if file_type == "pdf":
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text = "\n\n".join(page.get_text() for page in doc)
        doc.close()
        images = _pdf_bytes_to_images(file_bytes, max_slides, max_dim)

    elif file_type in ("pptx", "ppt"):
        # --- Text via python-pptx ---
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slide_texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
        text = "\n\n".join(slide_texts)

        # --- Images via LibreOffice → PDF → PyMuPDF ---
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, "brief.pptx")
            pdf_path = os.path.join(tmpdir, "brief.pdf")

            with open(pptx_path, "wb") as f:
                f.write(file_bytes)

            result = subprocess.run(
                [
                    "libreoffice", "--headless",
                    "--convert-to", "pdf",
                    "--outdir", tmpdir,
                    pptx_path,
                ],
                capture_output=True,
                timeout=90,
            )

            if result.returncode == 0 and os.path.exists(pdf_path):
                with open(pdf_path, "rb") as f:
                    pdf_bytes_converted = f.read()
                images = _pdf_bytes_to_images(pdf_bytes_converted, max_slides, max_dim)
            else:
                # LibreOffice not available or conversion failed — images unavailable
                stderr = result.stderr.decode(errors="replace") if result.stderr else ""
                print(f"[multimodal] LibreOffice conversion failed: {stderr[:300]}")

    return {
        "text": text.strip(),
        "images": images,
        "slide_count": len(images),
    }
